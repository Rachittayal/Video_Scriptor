from pptx import Presentation
from pptx.util import Inches
import os
import shutil

def create_ppt(transcript_file, images_folder, output_file):
    # Create a new presentation
    ppt = Presentation()
    
    # Add title slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = "Industrial Equipment Overview"
    slide.placeholders[1].text = "Generated Presentation"
    
    # Read transcript content
    if not os.path.exists(transcript_file):
        print(f"Transcript file '{transcript_file}' not found!")
        return

    with open(transcript_file, "r", encoding="utf-8") as file:
        content = file.read()
    
    paragraphs = content.split(".")  # Split transcript into sentences
    images = sorted(os.listdir(images_folder)) if os.path.exists(images_folder) else []

    # Add content slides
    for i, paragraph in enumerate(paragraphs):
        if paragraph.strip():
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
            slide.placeholders[1].text = paragraph.strip()
            
            # Add image to the slide (if available)
            if i < len(images):
                img_path = os.path.join(images_folder, images[i])
                if os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(4))
    
    # Save the presentation
    os.makedirs(os.path.dirname(output_file), exist_ok=True)  # Ensure output folder exists
    ppt.save(output_file)
    print(f"PPT saved to {output_file}")


def download_pptx(output_file, download_location):
    """
    Downloads the generated PPTX file to the specified download location.

    Args:
        output_file (str): Path to the generated PPTX file.
        download_location (str): Path where the PPTX file should be downloaded.
    """
    if not os.path.exists(output_file):
        print(f"Output file '{output_file}' not found!")
        return
    
    # Ensure the download location exists
    os.makedirs(download_location, exist_ok=True)
    
    # Define the destination path
    destination = os.path.join(download_location, os.path.basename(output_file))
    
    # Copy the PPTX file to the specified download location
    shutil.copy(output_file, destination)
    print(f"PPTX downloaded to {destination}")


if __name__ == "__main__":
    # Paths for input and output files
    transcript_file = "data/transcript.txt"
    images_folder = "data/extracted_images"
    output_file = "output/finall_presentation.pptx"
    download_location = r"C:\Users\rachi\OneDrive\Desktop\output_pptx" 

    # Create the PPT
    create_ppt(transcript_file, images_folder, output_file)
    
    # Download the PPTX to the specified location
    download_pptx(output_file, download_location)
